"""Document parser service — multi-format text extraction.

Supports: PDF (PyMuPDF), DOCX/DOC (python-docx), XLSX (openpyxl), XLS (xlrd),
PPTX (python-pptx), DXF/DWG (ezdxf + LibreDWG dwg2dxf), ZIP (stdlib),
CSV/JSON/XML/HTML (stdlib + optional bs4).
"""
from __future__ import annotations

import csv
import io
import json
import os
import subprocess
import tempfile
import zipfile
from collections import Counter
from pathlib import Path
from typing import Any, Dict

from fastapi import FastAPI, File, HTTPException, UploadFile

try:
    from common.runtime import ErrorBudget, setup_json_logger
    logger = setup_json_logger("document-parser", os.getenv("LOG_PATH", "/tmp/doc-parser.json.log"))
    _budget: Any = ErrorBudget(window_seconds=300)
except Exception:
    import logging
    logging.basicConfig(level=logging.INFO)
    logger = logging.getLogger("document-parser")
    _budget = None

app = FastAPI(title="Document Parser", version="1.0.0")

_MAX_ZIP_EXTRACTED = 50 * 1024 * 1024  # 50 MB total across all members
_MAX_ZIP_MEMBERS = 100

# ── optional heavy libs (graceful degradation) ───────────────────────
try:
    import fitz as _fitz
    _FITZ_OK = True
except ImportError:
    _fitz = None  # type: ignore[assignment]
    _FITZ_OK = False

try:
    import docx as _docx
    _DOCX_OK = True
except ImportError:
    _docx = None  # type: ignore[assignment]
    _DOCX_OK = False

try:
    import openpyxl as _openpyxl
    _OPENPYXL_OK = True
except ImportError:
    _openpyxl = None  # type: ignore[assignment]
    _OPENPYXL_OK = False

try:
    import xlrd as _xlrd
    _XLRD_OK = True
except ImportError:
    _xlrd = None  # type: ignore[assignment]
    _XLRD_OK = False

try:
    from pptx import Presentation as _Presentation
    _PPTX_OK = True
except ImportError:
    _Presentation = None  # type: ignore[assignment]
    _PPTX_OK = False

try:
    import ezdxf as _ezdxf
    _EZDXF_OK = True
except ImportError:
    _ezdxf = None  # type: ignore[assignment]
    _EZDXF_OK = False

try:
    from bs4 import BeautifulSoup as _BeautifulSoup
    _BS4_OK = True
except ImportError:
    _BeautifulSoup = None  # type: ignore[assignment]
    _BS4_OK = False


# ── parsers ──────────────────────────────────────────────────────────

def _parse_pdf(data: bytes, filename: str) -> Dict[str, Any]:
    if not _FITZ_OK:
        raise HTTPException(503, "PyMuPDF (fitz) not available — install pymupdf")
    doc = _fitz.open(stream=data, filetype="pdf")
    pages = [doc.load_page(i).get_text() for i in range(doc.page_count)]
    text = "\n\n".join(p.strip() for p in pages if p.strip())
    meta = doc.metadata or {}
    return {
        "text": text,
        "format": "pdf",
        "page_count": doc.page_count,
        "metadata": {"title": meta.get("title", ""), "author": meta.get("author", "")},
    }


def _parse_docx(data: bytes, filename: str) -> Dict[str, Any]:
    if not _DOCX_OK:
        raise HTTPException(503, "python-docx not available")
    doc = _docx.Document(io.BytesIO(data))
    parts: list[str] = []
    for para in doc.paragraphs:
        if para.text.strip():
            parts.append(para.text)
    for table in doc.tables:
        for row in table.rows:
            row_text = "\t".join(c.text.strip() for c in row.cells)
            if row_text.strip():
                parts.append(row_text)
    return {"text": "\n".join(parts), "format": "docx", "page_count": None, "metadata": {}}


def _parse_xlsx(data: bytes, filename: str) -> Dict[str, Any]:
    if not _OPENPYXL_OK:
        raise HTTPException(503, "openpyxl not available")
    wb = _openpyxl.load_workbook(io.BytesIO(data), read_only=True, data_only=True)
    sheet_names = list(wb.sheetnames)
    parts: list[str] = []
    for name in sheet_names:
        ws = wb[name]
        parts.append(f"=== Sheet: {name} ===")
        for row in ws.iter_rows(values_only=True):
            row_vals = [str(v) if v is not None else "" for v in row]
            if any(v.strip() for v in row_vals):
                parts.append("\t".join(row_vals))
    wb.close()
    return {
        "text": "\n".join(parts),
        "format": "xlsx",
        "page_count": len(sheet_names),
        "metadata": {"sheets": sheet_names},
    }


def _parse_xls(data: bytes, filename: str) -> Dict[str, Any]:
    if not _XLRD_OK:
        raise HTTPException(503, "xlrd not available")
    wb = _xlrd.open_workbook(file_contents=data)
    parts: list[str] = []
    for name in wb.sheet_names():
        ws = wb.sheet_by_name(name)
        parts.append(f"=== Sheet: {name} ===")
        for r in range(ws.nrows):
            row_vals = [str(ws.cell_value(r, c)) for c in range(ws.ncols)]
            if any(v.strip() for v in row_vals):
                parts.append("\t".join(row_vals))
    return {
        "text": "\n".join(parts),
        "format": "xls",
        "page_count": wb.nsheets,
        "metadata": {"sheets": wb.sheet_names()},
    }


def _parse_pptx(data: bytes, filename: str) -> Dict[str, Any]:
    if not _PPTX_OK:
        raise HTTPException(503, "python-pptx not available")
    prs = _Presentation(io.BytesIO(data))
    parts: list[str] = []
    for i, slide in enumerate(prs.slides, 1):
        slide_parts: list[str] = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    txt = para.text.strip()
                    if txt:
                        slide_parts.append(txt)
        if slide_parts:
            parts.append(f"--- Slide {i} ---\n" + "\n".join(slide_parts))
    return {
        "text": "\n\n".join(parts),
        "format": "pptx",
        "page_count": len(prs.slides),
        "metadata": {},
    }


def _parse_dxf_bytes(data: bytes, filename: str) -> Dict[str, Any]:
    if not _EZDXF_OK:
        raise HTTPException(503, "ezdxf not available")
    with tempfile.NamedTemporaryFile(suffix=".dxf", delete=False) as f:
        f.write(data)
        tmp = f.name
    try:
        doc, _ = _ezdxf.recover.readfile(tmp)
    except Exception as exc:
        raise HTTPException(400, f"Invalid DXF file: {exc}") from exc
    finally:
        os.unlink(tmp)

    layers = [layer.dxf.name for layer in doc.layers]
    entity_counts: Counter = Counter()
    annotations: list[str] = []

    msp = doc.modelspace()
    for entity in msp:
        entity_counts[entity.dxftype()] += 1
        dtype = entity.dxftype()
        if dtype in ("TEXT", "ATTDEF"):
            txt = entity.dxf.get("text", "").strip()
            if txt:
                annotations.append(txt)
        elif dtype == "MTEXT":
            raw = getattr(entity, "text", "") or ""
            txt = raw.strip()
            if txt:
                annotations.append(txt)

    blocks = [blk.name for blk in doc.blocks if not blk.name.startswith("*")]
    insunits = doc.header.get("$INSUNITS", None)

    parts = [f"Layers ({len(layers)}): {', '.join(layers[:100])}"]
    if entity_counts:
        parts.append(f"Entity types: {dict(entity_counts)}")
    if blocks:
        parts.append(f"Blocks ({len(blocks)}): {', '.join(blocks[:50])}")
    if annotations:
        parts.append(f"Annotations ({len(annotations)}):\n" + "\n".join(annotations[:300]))

    return {
        "text": "\n\n".join(parts),
        "format": "dxf",
        "page_count": None,
        "metadata": {
            "layers": layers[:100],
            "entity_counts": dict(entity_counts),
            "blocks": blocks[:50],
            "units": str(insunits),
        },
    }


def _parse_dwg(data: bytes, filename: str) -> Dict[str, Any]:
    if not _EZDXF_OK:
        raise HTTPException(503, "ezdxf not available")
    safe_name = Path(filename or "input.dwg").name
    with tempfile.TemporaryDirectory() as tmp:
        dwg_path = os.path.join(tmp, safe_name)
        dxf_path = dwg_path.rsplit(".", 1)[0] + ".dxf"
        with open(dwg_path, "wb") as f:
            f.write(data)
        try:
            result = subprocess.run(
                ["dwg2dxf", dwg_path, "-o", dxf_path],
                capture_output=True, timeout=30,
            )
        except FileNotFoundError:
            raise HTTPException(503, "dwg2dxf not available — install libredwg-tools")
        except subprocess.TimeoutExpired:
            raise HTTPException(504, "DWG conversion timed out")
        if result.returncode != 0:
            raise HTTPException(502, f"dwg2dxf error: {result.stderr.decode(errors='replace')[:300]}")
        with open(dxf_path, "rb") as f:
            dxf_data = f.read()
    parsed = _parse_dxf_bytes(dxf_data, safe_name.replace(".dwg", ".dxf"))
    parsed["format"] = "dwg"
    return parsed


def _parse_zip(data: bytes, filename: str) -> Dict[str, Any]:
    if not zipfile.is_zipfile(io.BytesIO(data)):
        raise HTTPException(400, "Not a valid ZIP file")
    parts: list[str] = []
    file_list: list[str] = []
    total_extracted = 0
    with zipfile.ZipFile(io.BytesIO(data)) as zf:
        for info in zf.infolist()[:_MAX_ZIP_MEMBERS]:
            if info.filename.endswith("/"):
                continue
            file_list.append(info.filename)
            if total_extracted >= _MAX_ZIP_EXTRACTED:
                parts.append(f"=== {info.filename} === [skipped — extraction limit reached]")
                continue
            try:
                member_data = zf.read(info.filename)
                total_extracted += len(member_data)
                ext = info.filename.rsplit(".", 1)[-1].lower() if "." in info.filename else ""
                parsed = _dispatch(member_data, info.filename, ext, depth=1)
                if parsed.get("text"):
                    parts.append(f"=== {info.filename} ===\n{parsed['text']}")
                else:
                    parts.append(f"=== {info.filename} === [empty]")
            except HTTPException:
                parts.append(f"=== {info.filename} === [unsupported format]")
            except Exception as exc:
                parts.append(f"=== {info.filename} === [error: {exc}]")
    return {
        "text": "\n\n".join(parts),
        "format": "zip",
        "page_count": len(file_list),
        "metadata": {"files": file_list},
    }


def _parse_csv(data: bytes, filename: str) -> Dict[str, Any]:
    text = data.decode("utf-8", errors="replace")
    reader = csv.reader(io.StringIO(text))
    rows = list(reader)
    formatted = "\n".join("\t".join(row) for row in rows)
    return {
        "text": formatted,
        "format": "csv",
        "page_count": None,
        "metadata": {"rows": len(rows)},
    }


def _parse_xml_html(data: bytes, filename: str, fmt: str) -> Dict[str, Any]:
    import re
    text_raw = data.decode("utf-8", errors="replace")
    if _BS4_OK:
        parser = "lxml" if fmt == "xml" else "html.parser"
        soup = _BeautifulSoup(text_raw, parser)
        text = soup.get_text(separator="\n", strip=True)
    else:
        text = re.sub(r"<[^>]+>", "", text_raw)
    return {"text": text, "format": fmt, "page_count": None, "metadata": {}}


def _parse_json(data: bytes, filename: str) -> Dict[str, Any]:
    try:
        obj = json.loads(data.decode("utf-8", errors="replace"))
        text = json.dumps(obj, indent=2, ensure_ascii=False)
    except Exception:
        text = data.decode("utf-8", errors="replace")
    return {"text": text, "format": "json", "page_count": None, "metadata": {}}


# ── dispatch ─────────────────────────────────────────────────────────

def _dispatch(data: bytes, filename: str, ext: str, depth: int = 0) -> Dict[str, Any]:
    if ext == "pdf":
        return _parse_pdf(data, filename)
    if ext in ("docx", "doc"):
        return _parse_docx(data, filename)
    if ext == "xlsx":
        return _parse_xlsx(data, filename)
    if ext == "xls":
        return _parse_xls(data, filename)
    if ext in ("pptx", "ppt"):
        return _parse_pptx(data, filename)
    if ext == "dxf":
        return _parse_dxf_bytes(data, filename)
    if ext == "dwg":
        if depth > 0:
            return {"text": "[DWG inside ZIP not supported]", "format": "dwg", "page_count": None, "metadata": {}}
        return _parse_dwg(data, filename)
    if ext == "zip" and depth == 0:
        return _parse_zip(data, filename)
    if ext == "csv":
        return _parse_csv(data, filename)
    if ext in ("xml", "html", "htm"):
        return _parse_xml_html(data, filename, ext)
    if ext == "json":
        return _parse_json(data, filename)
    # Plain-text fallback
    return {"text": data.decode("utf-8", errors="replace"), "format": ext or "text", "page_count": None, "metadata": {}}


# ── endpoints ────────────────────────────────────────────────────────

@app.post("/parse")
async def parse(file: UploadFile = File(...)) -> Dict[str, Any]:
    """Parse an uploaded document and return extracted text."""
    if not file.filename:
        raise HTTPException(400, "No filename provided")
    data = await file.read()
    if not data:
        raise HTTPException(400, "Empty file")
    ext = (file.filename.rsplit(".", 1)[-1] if "." in file.filename else "").lower()
    result = _dispatch(data, file.filename, ext)
    result["filename"] = file.filename
    logger.info("Parsed %s (%s) → %d chars", file.filename, result["format"], len(result.get("text", "")))
    return result


@app.get("/formats")
async def formats() -> Dict[str, Any]:
    """List supported formats and library availability."""
    _dwg_tool = subprocess.run(["which", "dwg2dxf"], capture_output=True).returncode == 0
    return {
        "pdf": _FITZ_OK,
        "docx": _DOCX_OK,
        "doc": _DOCX_OK,
        "xlsx": _OPENPYXL_OK,
        "xls": _XLRD_OK,
        "pptx": _PPTX_OK,
        "ppt": _PPTX_OK,
        "dxf": _EZDXF_OK,
        "dwg": _EZDXF_OK and _dwg_tool,
        "zip": True,
        "csv": True,
        "json": True,
        "xml": True,
        "html": True,
    }


@app.get("/health")
async def health() -> Dict[str, Any]:
    return {
        "status": "ok",
        "pdf": _FITZ_OK,
        "docx": _DOCX_OK,
        "xlsx": _OPENPYXL_OK,
        "xls": _XLRD_OK,
        "pptx": _PPTX_OK,
        "dxf": _EZDXF_OK,
    }


@app.get("/metrics")
async def metrics() -> Dict[str, Any]:
    snap = _budget.snapshot() if _budget else {}
    return {"status": "ok", "error_budget": snap}


if __name__ == "__main__":
    import uvicorn
    uvicorn.run(app, host="0.0.0.0", port=8032)
